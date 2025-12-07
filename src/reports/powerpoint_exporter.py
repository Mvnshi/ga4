"""
PowerPoint Exporter
===================
Professional presentation generation for executive reports.

Creates branded slide decks with:
- Executive summary
- Key metrics with YoY comparison
- Visual charts and graphs
- Insights and recommendations
"""

from pathlib import Path
from typing import Dict, Any, List, Tuple
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

from config.settings import OUTPUT_DIR, ClientConfig


class PowerPointExporter:
    """
    Creates professional PowerPoint presentations.
    
    Generates McKinsey-style slide decks with:
    - Title slide
    - Executive summary
    - Key metrics dashboard
    - Detailed analysis slides
    - Recommendations
    """
    
    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Colors (RGB tuples)
    COLORS = {
        'primary': RgbColor(45, 80, 22),      # Forest green
        'secondary': RgbColor(244, 196, 48),   # Honey gold
        'dark': RgbColor(31, 41, 55),          # Dark gray
        'light': RgbColor(248, 250, 252),      # Light gray
        'positive': RgbColor(34, 197, 94),     # Green
        'negative': RgbColor(239, 68, 68),     # Red
        'white': RgbColor(255, 255, 255),
        'black': RgbColor(0, 0, 0),
    }
    
    # Margins
    MARGIN = Inches(0.5)
    
    def __init__(self, client_config: ClientConfig):
        """Initialize exporter with client configuration."""
        self.config = client_config
        self.prs = None
    
    def export(
        self,
        report_data: Dict[str, Any],
        filename: str = None
    ) -> Path:
        """
        Export report data to PowerPoint.
        
        Args:
            report_data: Full report data dictionary
            filename: Output filename (auto-generated if not provided)
        
        Returns:
            Path to created PowerPoint file
        """
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        
        # Create slides
        self._create_title_slide(report_data)
        self._create_executive_summary(report_data)
        self._create_traffic_overview(report_data)
        self._create_traffic_trends(report_data)
        self._create_search_performance(report_data)
        self._create_top_keywords(report_data)
        self._create_content_performance(report_data)
        self._create_audience_breakdown(report_data)
        self._create_acquisition_channels(report_data)
        self._create_insights_slide(report_data)
        self._create_recommendations_slide(report_data)
        self._create_closing_slide(report_data)
        
        # Generate filename
        if filename is None:
            period = report_data.get('metadata', {}).get('current_period', {}).get('label', 'report')
            client_name = self.config.name
            filename = f"{client_name}_quarterly_presentation_{period.replace(' ', '_')}.pptx"
        
        output_path = OUTPUT_DIR / filename
        self.prs.save(output_path)
        
        return output_path
    
    def _add_blank_slide(self) -> object:
        """Add a blank slide."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(blank_layout)
    
    def _add_title(
        self,
        slide,
        text: str,
        top: float = 0.3,
        font_size: int = 32
    ):
        """Add a title to a slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(top),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
    
    def _add_subtitle(
        self,
        slide,
        text: str,
        top: float = 0.9
    ):
        """Add a subtitle to a slide."""
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(top),
            Inches(12.333), Inches(0.5)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = self.COLORS['dark']
    
    def _add_metric_card(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        label: str,
        value: str,
        change: str = None,
        is_positive: bool = True
    ):
        """Add a metric card with value and change indicator."""
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['light']
        shape.line.fill.background()
        
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.1),
            Inches(width - 0.3), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = self.COLORS['dark']
        
        # Value
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.45),
            Inches(width - 0.3), Inches(0.5)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # Change indicator
        if change:
            change_box = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(top + 0.95),
                Inches(width - 0.3), Inches(0.3)
            )
            tf = change_box.text_frame
            p = tf.paragraphs[0]
            arrow = "↑" if is_positive else "↓"
            p.text = f"{arrow} {change}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['positive'] if is_positive else self.COLORS['negative']
    
    def _add_table(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        data: List[List[str]],
        headers: List[str],
        col_widths: List[float] = None
    ):
        """Add a formatted table to the slide."""
        rows = len(data) + 1  # +1 for header
        cols = len(headers)
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(0.4 * rows)
        ).table
        
        # Set column widths
        if col_widths:
            for i, cw in enumerate(col_widths):
                table.columns[i].width = Inches(cw)
        
        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['primary']
            
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = self.COLORS['white']
            para.alignment = PP_ALIGN.CENTER
        
        # Data rows
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(9)
                para.font.color.rgb = self.COLORS['dark']
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['light']
    
    def _create_title_slide(self, data: Dict[str, Any]):
        """Create title slide."""
        slide = self._add_blank_slide()
        
        # Background accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), self.SLIDE_HEIGHT
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.COLORS['primary']
        accent.line.fill.background()
        
        # Organization name
        org_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11), Inches(1)
        )
        tf = org_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.config.display_name
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.2),
            Inches(11), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Quarterly Analytics Report"
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLORS['dark']
        
        # Period
        metadata = data.get('metadata', {})
        current = metadata.get('current_period', {})
        period = current.get('label', 'Current Quarter')
        
        period_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2),
            Inches(11), Inches(0.5)
        )
        tf = period_box.text_frame
        p = tf.paragraphs[0]
        p.text = period
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['secondary']
        p.font.bold = True
    
    def _create_executive_summary(self, data: Dict[str, Any]):
        """Create executive summary slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Executive Summary")
        
        insights = data.get('insights', {})
        summary = insights.get('executive_summary', 'Analysis pending.')
        
        # Summary box
        summary_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(12.333), Inches(1.5)
        )
        tf = summary_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = summary
        p.font.size = Pt(14)
        p.font.color.rgb = self.COLORS['dark']
        
        # Key findings section
        findings_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(3),
            Inches(12.333), Inches(0.5)
        )
        tf = findings_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Findings"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # List insights
        insights_list = insights.get('insights', [])
        findings_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5),
            Inches(12.333), Inches(3.5)
        )
        tf = findings_box.text_frame
        tf.word_wrap = True
        
        for i, insight in enumerate(insights_list[:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            emoji = "✓" if insight.get('type') == 'positive' else "⚠" if insight.get('type') == 'negative' else "→"
            p.text = f"{emoji}  {insight.get('headline', '')}"
            p.font.size = Pt(12)
            p.space_after = Pt(8)
            
            # Color by type
            if insight.get('type') == 'positive':
                p.font.color.rgb = self.COLORS['positive']
            elif insight.get('type') == 'negative':
                p.font.color.rgb = self.COLORS['negative']
            else:
                p.font.color.rgb = self.COLORS['dark']
    
    def _create_traffic_overview(self, data: Dict[str, Any]):
        """Create traffic overview slide with metric cards."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Website Traffic Overview")
        
        metadata = data.get('metadata', {})
        current = metadata.get('current_period', {})
        previous = metadata.get('previous_period', {})
        self._add_subtitle(slide, f"{current.get('label', '')} vs {previous.get('label', '')}")
        
        traffic = data.get('ga4', {}).get('traffic_overview', {})
        comparison = data.get('comparison', {}).get('traffic_overview', {})
        
        # Metric cards in 2x3 grid
        metrics = [
            ('Total Users', 'total_users'),
            ('New Users', 'new_users'),
            ('Sessions', 'sessions'),
            ('Pageviews', 'pageviews'),
            ('Bounce Rate', 'bounce_rate'),
            ('Avg Session (sec)', 'avg_session_duration'),
        ]
        
        card_width = 3.8
        card_height = 1.4
        start_left = 0.7
        start_top = 1.6
        gap = 0.3
        
        for i, (label, key) in enumerate(metrics):
            row = i // 3
            col = i % 3
            
            left = start_left + col * (card_width + gap)
            top = start_top + row * (card_height + gap)
            
            value = traffic.get(key, 0)
            comp_data = comparison.get(key, {})
            
            if isinstance(value, float) and key in ['bounce_rate', 'avg_session_duration']:
                value_str = f"{value:.1f}"
            else:
                value_str = f"{value:,}"
            
            change_data = comp_data.get('change', {}) if isinstance(comp_data, dict) else {}
            change_str = change_data.get('formatted', '') if isinstance(change_data, dict) else ''
            direction = change_data.get('direction', 'neutral') if isinstance(change_data, dict) else 'neutral'
            
            # For bounce rate, down is good
            is_positive = direction == 'up'
            if key == 'bounce_rate':
                is_positive = direction == 'down'
            
            self._add_metric_card(
                slide, left, top, card_width, card_height,
                label, value_str, change_str, is_positive
            )
    
    def _create_traffic_trends(self, data: Dict[str, Any]):
        """Create traffic trends slide with monthly data."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Monthly Traffic Trends")
        
        monthly = data.get('ga4', {}).get('traffic_by_month', pd.DataFrame())
        
        if monthly.empty:
            return
        
        # Table data
        headers = ['Month', 'Users', 'New', 'Returning', 'Sessions', 'Bounce Rate']
        table_data = []
        
        for _, row in monthly.iterrows():
            table_data.append([
                str(row.get('yearMonth', '')),
                f"{int(row.get('totalUsers', 0)):,}",
                f"{int(row.get('newUsers', 0)):,}",
                f"{int(row.get('returning_users', 0)):,}",
                f"{int(row.get('sessions', 0)):,}",
                f"{row.get('bounceRate', 0):.1f}%",
            ])
        
        self._add_table(
            slide, 0.5, 1.5, 12.333,
            table_data[-6:],  # Last 6 months
            headers,
            col_widths=[2, 2, 2, 2, 2, 2]
        )
    
    def _create_search_performance(self, data: Dict[str, Any]):
        """Create search performance overview slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Search Performance Overview")
        
        gsc = data.get('gsc', {})
        overview = gsc.get('overview', {})
        prev_overview = data.get('comparison', {}).get('gsc', {}).get('overview', {})
        
        # Metric cards
        metrics = [
            ('Total Clicks', overview.get('total_clicks', 0), prev_overview.get('total_clicks', 0)),
            ('Total Impressions', overview.get('total_impressions', 0), prev_overview.get('total_impressions', 0)),
            ('Average CTR', f"{overview.get('avg_ctr', 0):.2f}%", None),
            ('Average Position', f"{overview.get('avg_position', 0):.1f}", None),
        ]
        
        card_width = 2.9
        card_height = 1.3
        start_left = 0.7
        
        for i, (label, value, prev) in enumerate(metrics):
            left = start_left + i * (card_width + 0.2)
            
            if prev is not None and prev > 0:
                curr_num = value if isinstance(value, (int, float)) else 0
                change_pct = ((curr_num - prev) / prev) * 100
                change_str = f"{'+' if change_pct > 0 else ''}{change_pct:.1f}%"
                is_positive = change_pct > 0
            else:
                change_str = None
                is_positive = True
            
            value_str = f"{value:,}" if isinstance(value, int) else str(value)
            
            self._add_metric_card(
                slide, left, 1.5, card_width, card_height,
                label, value_str, change_str, is_positive
            )
    
    def _create_top_keywords(self, data: Dict[str, Any]):
        """Create top keywords slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Top Search Keywords")
        self._add_subtitle(slide, "Keywords driving organic traffic")
        
        keywords = data.get('gsc', {}).get('top_keywords_clicks', pd.DataFrame())
        
        if keywords.empty:
            return
        
        headers = ['Keyword', 'Clicks', 'Impressions', 'CTR', 'Position']
        table_data = []
        
        for _, row in keywords.head(12).iterrows():
            # Truncate long keywords
            keyword = str(row.get('query', ''))[:40]
            if len(str(row.get('query', ''))) > 40:
                keyword += '...'
            
            table_data.append([
                keyword,
                f"{int(row.get('clicks', 0)):,}",
                f"{int(row.get('impressions', 0)):,}",
                f"{row.get('ctr', 0):.2f}%",
                f"{row.get('position', 0):.1f}",
            ])
        
        self._add_table(
            slide, 0.5, 1.6, 12.333,
            table_data,
            headers,
            col_widths=[5, 1.8, 2, 1.5, 1.5]
        )
    
    def _create_content_performance(self, data: Dict[str, Any]):
        """Create top content slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Top Performing Content")
        
        top_pages = data.get('ga4', {}).get('top_pages', pd.DataFrame())
        
        if top_pages.empty:
            return
        
        headers = ['Page', 'Views', '% Total', 'Avg Time', 'Bounce']
        table_data = []
        
        for _, row in top_pages.head(10).iterrows():
            title = str(row.get('pageTitle', row.get('pagePath', '')))[:45]
            if len(str(row.get('pageTitle', ''))) > 45:
                title += '...'
            
            table_data.append([
                title,
                f"{int(row.get('screenPageViews', 0)):,}",
                f"{row.get('pct_of_total', 0):.1f}%",
                f"{row.get('averageSessionDuration', 0):.0f}s",
                f"{row.get('bounceRate', 0):.1f}%",
            ])
        
        self._add_table(
            slide, 0.5, 1.4, 12.333,
            table_data,
            headers,
            col_widths=[5.5, 1.5, 1.5, 1.5, 1.5]
        )
    
    def _create_audience_breakdown(self, data: Dict[str, Any]):
        """Create audience insights slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Audience Breakdown")
        
        # Device data
        devices = data.get('ga4', {}).get('device_breakdown', pd.DataFrame())
        
        if not devices.empty:
            device_title = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3),
                Inches(6), Inches(0.4)
            )
            tf = device_title.text_frame
            p = tf.paragraphs[0]
            p.text = "By Device"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            
            headers = ['Device', 'Users', 'Share', 'Bounce Rate']
            table_data = []
            
            for _, row in devices.iterrows():
                table_data.append([
                    str(row.get('deviceCategory', '')).title(),
                    f"{int(row.get('totalUsers', 0)):,}",
                    f"{row.get('user_share', 0):.1f}%",
                    f"{row.get('bounceRate', 0):.1f}%",
                ])
            
            self._add_table(
                slide, 0.5, 1.8, 5.5,
                table_data,
                headers,
                col_widths=[1.5, 1.5, 1.25, 1.25]
            )
        
        # Geography
        geo = data.get('ga4', {}).get('geography', pd.DataFrame())
        
        if not geo.empty:
            geo_title = slide.shapes.add_textbox(
                Inches(7), Inches(1.3),
                Inches(6), Inches(0.4)
            )
            tf = geo_title.text_frame
            p = tf.paragraphs[0]
            p.text = "Top Countries"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            
            headers = ['Country', 'Users', 'Share']
            table_data = []
            
            for _, row in geo.head(8).iterrows():
                table_data.append([
                    str(row.get('country', ''))[:20],
                    f"{int(row.get('totalUsers', 0)):,}",
                    f"{row.get('user_share', 0):.1f}%",
                ])
            
            self._add_table(
                slide, 7, 1.8, 5.5,
                table_data,
                headers,
                col_widths=[2.5, 1.5, 1.5]
            )
    
    def _create_acquisition_channels(self, data: Dict[str, Any]):
        """Create acquisition channels slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Traffic Acquisition Channels")
        
        channels = data.get('ga4', {}).get('traffic_by_channel', pd.DataFrame())
        
        if channels.empty:
            return
        
        headers = ['Channel', 'Sessions', 'Share', 'Bounce Rate', 'Avg Duration']
        table_data = []
        
        for _, row in channels.head(8).iterrows():
            table_data.append([
                str(row.get('sessionDefaultChannelGroup', '')),
                f"{int(row.get('sessions', 0)):,}",
                f"{row.get('session_share', 0):.1f}%",
                f"{row.get('bounceRate', 0):.1f}%",
                f"{row.get('averageSessionDuration', 0):.0f}s",
            ])
        
        self._add_table(
            slide, 0.5, 1.5, 12.333,
            table_data,
            headers,
            col_widths=[3.5, 2, 2, 2.2, 2.2]
        )
    
    def _create_insights_slide(self, data: Dict[str, Any]):
        """Create key insights slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Key Insights")
        
        insights = data.get('insights', {}).get('insights', [])
        
        if not insights:
            return
        
        # Group by type
        positive = [i for i in insights if i.get('type') == 'positive'][:2]
        negative = [i for i in insights if i.get('type') == 'negative'][:2]
        opportunities = [i for i in insights if i.get('type') == 'opportunity'][:2]
        
        current_top = 1.4
        
        for section_name, section_insights, color in [
            ("Strengths", positive, self.COLORS['positive']),
            ("Challenges", negative, self.COLORS['negative']),
            ("Opportunities", opportunities, self.COLORS['secondary']),
        ]:
            if section_insights:
                # Section header
                header_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(current_top),
                    Inches(12.333), Inches(0.4)
                )
                tf = header_box.text_frame
                p = tf.paragraphs[0]
                p.text = section_name
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = color
                
                current_top += 0.4
                
                # Insights
                for insight in section_insights:
                    insight_box = slide.shapes.add_textbox(
                        Inches(0.7), Inches(current_top),
                        Inches(12.133), Inches(0.5)
                    )
                    tf = insight_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"• {insight.get('headline', '')}"
                    p.font.size = Pt(11)
                    p.font.color.rgb = self.COLORS['dark']
                    
                    current_top += 0.45
                
                current_top += 0.3
    
    def _create_recommendations_slide(self, data: Dict[str, Any]):
        """Create recommendations slide."""
        slide = self._add_blank_slide()
        self._add_title(slide, "Recommendations")
        
        recommendations = data.get('insights', {}).get('key_recommendations', [])
        
        if not recommendations:
            return
        
        current_top = 1.5
        
        for i, rec in enumerate(recommendations[:5], 1):
            # Number box
            num_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(current_top),
                Inches(0.4), Inches(0.4)
            )
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = self.COLORS['primary']
            num_shape.line.fill.background()
            
            num_tf = num_shape.text_frame
            num_tf.paragraphs[0].text = str(i)
            num_tf.paragraphs[0].font.color.rgb = self.COLORS['white']
            num_tf.paragraphs[0].font.bold = True
            num_tf.paragraphs[0].font.size = Pt(14)
            num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Recommendation text
            rec_box = slide.shapes.add_textbox(
                Inches(1.1), Inches(current_top),
                Inches(11.5), Inches(0.8)
            )
            tf = rec_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = rec
            p.font.size = Pt(12)
            p.font.color.rgb = self.COLORS['dark']
            
            current_top += 1.0
    
    def _create_closing_slide(self, data: Dict[str, Any]):
        """Create closing/thank you slide."""
        slide = self._add_blank_slide()
        
        # Background accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.5),
            self.SLIDE_WIDTH, Inches(1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.COLORS['primary']
        accent.line.fill.background()
        
        # Thank you text
        thank_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1)
        )
        tf = thank_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.6),
            Inches(12.333), Inches(0.5)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Questions & Discussion"
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

